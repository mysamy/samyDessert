from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import re
import os
import io
from pygments import highlight
from pygments.lexers import PhpLexer, JavascriptLexer, HtmlDjangoLexer
from pygments.formatters import ImageFormatter
from pygments.style import Style
from pygments.token import (Token, Comment, Keyword, Name,
                             String, Number, Operator, Punctuation)

class VscodeDarkPlus(Style):
    background_color = '#1E1E1E'
    default_style = ''
    styles = {
        Token:               '#D4D4D4',
        Comment:             '#6A9955',
        Keyword:             '#569CD6',
        Keyword.Type:        '#4EC9B0',
        Keyword.Declaration: '#569CD6',
        Name.Function:       '#DCDCAA',
        Name.Class:          '#4EC9B0',
        Name.Builtin:        '#DCDCAA',
        Name.Variable:       '#9CDCFE',
        Name.Attribute:      '#9CDCFE',
        Name.Tag:            '#569CD6',
        String:              '#CE9178',
        String.Interpol:     '#CE9178',
        Number:              '#B5CEA8',
        Operator:            '#D4D4D4',
        Punctuation:         '#D4D4D4',
    }

CAPTURES = r"c:\Travail\samyDessert\Soutenance\DossierProjet\captures\\"

# Couleurs
CHOCOLAT   = RGBColor(0x3B, 0x1F, 0x0E)
FRAMBOISE  = RGBColor(0xAD, 0x14, 0x57)
PISTACHE   = RGBColor(0x7B, 0xAE, 0x7F)
CREME      = RGBColor(0xFD, 0xF6, 0xEC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG    = RGBColor(0x1E, 0x1E, 0x1E)
CODE_FG    = RGBColor(0xD4, 0xD4, 0xD4)
GREY_LIGHT = RGBColor(0xE8, 0xE0, 0xD8)
GREY_MID   = RGBColor(0xC8, 0xBE, 0xB4)

W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CREME
    return slide

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h, size=18, bold=False, color=CHOCOLAT,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def add_title_bar(slide, title, subtitle=None):
    # Barre framboise en haut
    add_rect(slide, 0, 0, W, Inches(1.2), fill=FRAMBOISE)
    add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9),
                size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.4),
                    size=13, bold=False, color=FRAMBOISE, italic=True)

def add_bullet_text(slide, lines, x, y, w, h, size=19):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        # Détecter niveau d'indentation
        indent = 0
        stripped = line.lstrip()
        if line.startswith('  ') or line.startswith('\t'):
            indent = 1
        # Détecter bullet
        is_bullet = stripped.startswith('- ') or stripped.startswith('• ')
        if is_bullet:
            text = stripped[2:]
            p.level = indent
        else:
            text = stripped
            p.level = 0
        # Parser le gras/italique basique
        parts = re.split(r'(\*\*.*?\*\*|`[^`]+`)', text)
        for part in parts:
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            elif part.startswith('`') and part.endswith('`'):
                run.text = part[1:-1]
                run.font.name = "Courier New"
                run.font.color.rgb = FRAMBOISE
            else:
                run.text = part
            run.font.size = Pt(size)
            run.font.color.rgb = CHOCOLAT
            run.font.name = run.font.name or "Calibri"
    return tb

def add_code_block(slide, code, x, y, w, h, lang='php'):
    lexers = {
        'php':  PhpLexer(startinline=True),
        'js':   JavascriptLexer(),
        'twig': HtmlDjangoLexer(),
    }
    lexer = lexers.get(lang, PhpLexer(startinline=True))
    fmt = ImageFormatter(
        style=VscodeDarkPlus,
        font_size=22,
        line_numbers=False,
        image_pad=12,
    )
    png = highlight(code, lexer, fmt)
    slide.shapes.add_picture(io.BytesIO(png), x, y, width=w)

def add_table(slide, headers, rows, x, y, w):
    cols = len(headers)
    n_rows = len(rows) + 1
    col_w = w // cols
    row_h = Inches(0.42)
    tbl = slide.shapes.add_table(n_rows, cols, x, y, w, row_h * n_rows).table
    tbl.columns[0].width = col_w

    # Header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = FRAMBOISE
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.size = Pt(15)
        p.alignment = PP_ALIGN.CENTER

    # Rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREY_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            # Parse backticks
            parts = re.split(r'(`[^`]+`)', val)
            for part in parts:
                run = p.add_run()
                if part.startswith('`') and part.endswith('`'):
                    run.text = part[1:-1]
                    run.font.name = "Courier New"
                    run.font.color.rgb = FRAMBOISE
                else:
                    run.text = part
                run.font.size = Pt(14)
                run.font.color.rgb = CHOCOLAT

def add_image_placeholder(slide, label, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=GREY_LIGHT, line=GREY_MID)
    add_textbox(slide, f"[ {label} ]", x, y + h//2 - Inches(0.2), w, Inches(0.4),
                size=11, color=GREY_MID, align=PP_ALIGN.CENTER, italic=True)

def add_image(slide, path, x, y, w, h):
    if not os.path.exists(path):
        add_image_placeholder(slide, os.path.basename(path), x, y, w, h)
        return
    slide.shapes.add_picture(path, x, y, w, h)

# ─── GÉNÉRATION DES SLIDES ──────────────────────────────────────────────────

prs = new_prs()

# ── SLIDE 1 — Titre ──────────────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, W, H, fill=CHOCOLAT)
add_rect(s, 0, Inches(2.5), W, Inches(2.8), fill=FRAMBOISE)
add_textbox(s, "SamyDessert", Inches(1), Inches(2.6), Inches(11), Inches(1.2),
            size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Boutique en ligne de pâtisseries artisanales",
            Inches(1), Inches(3.7), Inches(11), Inches(0.6),
            size=22, color=CREME, align=PP_ALIGN.CENTER)
add_textbox(s, "Samy Ben Hamida  ·  Développeur Web DWWM  ·  ESRP Auxilia  ·  2026",
            Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            size=14, color=GREY_LIGHT, align=PP_ALIGN.CENTER)

# ── SLIDE 2 — Sommaire ───────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Sommaire")
items = [
    "1. Contexte & Analyse des besoins",
    "2. Design & Identité visuelle",
    "3. Stack & Architecture technique",
    "4. Fonctionnalités",
    "5. Qualité & Sécurité",
    "6. Déploiement & Organisation",
    "7. Bilan & Perspectives",
]
add_bullet_text(s, items, Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), size=24)

# ── SLIDE 3 — Présentation personnelle ───────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Présentation personnelle & projet")
add_bullet_text(s, [
    "**Samy Ben Hamida — 34 ans**",
    "- Formation **Développeur Web et Web Mobile** (RNCP Niveau 5) — ESRP Auxilia, alternance",
    "- Stage chez **Creative Handicap** — association qui rend le numérique accessible à tous",
    "",
    "**SamyDessert — pourquoi ce projet ?**",
    "- Passionné de pâtisserie, je partage régulièrement mes desserts avec mes collègues",
    "- Concept hybride : recettes gratuites + boutique en ligne",
], Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.5), size=19)
add_image(s, CAPTURES + "accueil.png", Inches(8.2), Inches(1.4), Inches(4.7), Inches(4.5))

# ── SLIDE 4 — Benchmark ───────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Benchmark & Inspiration")
add_bullet_text(s, [
    "**Analyse de sites concurrents avant de commencer**",
    "",
    "Points retenus :",
    "- Mise en avant visuelle des produits (grandes photos)",
    "- Navigation claire par catégorie",
    "- Identité visuelle forte et cohérente",
], Inches(0.5), Inches(1.4), Inches(7.5), Inches(4), size=19)
add_image(s, CAPTURES + "Inspiration.png", Inches(8.2), Inches(1.4), Inches(4.7), Inches(4.5))

# ── SLIDE 5 — Personas ───────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Personas — Qui sont les utilisateurs cibles ?")
add_image(s, CAPTURES + "personna1.png", Inches(0.5), Inches(1.5), Inches(6), Inches(4.5))
add_image(s, CAPTURES + "personna2.png", Inches(6.9), Inches(1.5), Inches(6), Inches(4.5))

# ── SLIDE 6 — User Flow ──────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "User Flow — Parcours utilisateur principal")
add_textbox(s,
    "Accueil  →  Catalogue  →  Fiche produit  →  Panier  →  Commande  →  Paiement  →  Email de confirmation",
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
    size=16, bold=True, color=FRAMBOISE, align=PP_ALIGN.CENTER)
add_image(s, CAPTURES + "userFlow.png", Inches(0.5), Inches(2.3), Inches(12.3), Inches(4))

# ── SLIDE 7 — Wireframe ──────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Wireframe — Maquettes basse fidélité")
add_bullet_text(s, [
    "Permet de valider la structure des pages avant de coder :",
    "- Disposition des éléments",
    "- Navigation",
    "- Hiérarchie de l'information",
], Inches(0.5), Inches(1.4), Inches(5.5), Inches(2.5), size=19)
add_image(s, CAPTURES + "wireframeLow.png", Inches(6.2), Inches(1.4), Inches(6.7), Inches(5.5))

# ── SLIDE 8 — Identité visuelle ──────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Identité visuelle & Design system")
# Palette
y0 = Inches(1.4)
colors_info = [
    (RGBColor(0x3B,0x1F,0x0E), "Chocolat #3B1F0E", "Textes, fonds sombres", WHITE),
    (RGBColor(0xC0,0x39,0x2B), "Framboise #C0392B", "CTA, accents", WHITE),
    (RGBColor(0x7B,0xAE,0x7F), "Pistache #7BAE7F", "Badges, succès", WHITE),
    (RGBColor(0xFD,0xF6,0xEC), "Crème #FDF6EC", "Fond principal", CHOCOLAT),
]
for i, (col, name, role, txt_col) in enumerate(colors_info):
    xi = Inches(0.4 + i * 2.15)
    add_rect(s, xi, y0, Inches(1.9), Inches(1.1), fill=col)
    add_textbox(s, name, xi, y0 + Inches(0.05), Inches(1.9), Inches(0.5),
                size=10, bold=True, color=txt_col, align=PP_ALIGN.CENTER)
    add_textbox(s, role, xi, y0 + Inches(0.55), Inches(1.9), Inches(0.4),
                size=9, color=txt_col, align=PP_ALIGN.CENTER)

# Composants
add_bullet_text(s, [
    "**Composants réutilisables — Atomic Design :**",
    "- **Atoms** : bouton, badge, input",
    "- **Molecules** : input + label",
    "- **Organisms** : header, footer",
], Inches(0.4), Inches(2.7), Inches(5), Inches(2.5), size=20)
add_textbox(s, "Logo — \"S\" stylisé, conçu avec Affinity Designer",
            Inches(5.5), Inches(2.7), Inches(7.4), Inches(0.4),
            size=13, bold=True, color=CHOCOLAT)
add_image(s, CAPTURES + "Logo.png", Inches(5.5), Inches(3.15), Inches(2.3), Inches(1.7))
add_image(s, CAPTURES + "token.png", Inches(8.2), Inches(2.7), Inches(2.5), Inches(2))
add_image(s, CAPTURES + "Button.png", Inches(10.9), Inches(2.7), Inches(1.0), Inches(2))
add_image(s, CAPTURES + "InputField.png", Inches(11.9), Inches(2.7), Inches(1.0), Inches(2))

# ── SLIDE 9 — Responsive & Accessibilité ────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Responsive Design & Accessibilité")
add_bullet_text(s, [
    "**Mobile-first avec Tailwind CSS v4** — 3 breakpoints : Mobile / Tablette / Desktop",
    "",
    "**Accessibilité :**",
    "- Contraste vérifié (Colorable) — conformité WCAG",
    "- Police Luciole (dyslexie)",
    "- Focus visible sur tous les éléments interactifs",
], Inches(0.5), Inches(1.4), Inches(6.5), Inches(3), size=18)
for _i, _fname in enumerate(["grilleProduitMobile.png", "grilleProduitTablette.png",
                              "grilleProduitDesktop.png", "contrastColorable.png"]):
    add_image(s, CAPTURES + _fname, Inches(0.5 + _i * 3.1), Inches(4.5), Inches(2.9), Inches(2.5))

# ── SLIDE 10 — Stack technique ───────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Stack technique")
add_table(s,
    ["Couche", "Technologie"],
    [
        ["Backend",          "PHP 8.3 + Symfony 7.4"],
        ["Frontend",         "Twig + Tailwind CSS v4 + Stimulus"],
        ["Temps réel",       "Symfony UX Live Components (Turbo)"],
        ["Base de données",  "MySQL 8.0 + Doctrine ORM"],
        ["Paiement",         "Stripe (Checkout + Webhooks)"],
        ["Emails",           "Resend API"],
        ["Déploiement",      "Docker + Railway"],
        ["Assets",           "AssetMapper + ImportMap"],
    ],
    Inches(1), Inches(1.4), Inches(11.3)
)

# ── SLIDE 11 — Architecture MVC ──────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Architecture MVC")
add_bullet_text(s, [
    "**Models** : entités Doctrine (`Produit`, `Commande`, `Utilisateur`...)",
    "**Views** : templates Twig + composants atomiques",
    "**Controllers** : logique légère, délèguent aux Services",
    "**Services** : logique métier (`PanierService`, `MailerService`, `FactureService`)",
], Inches(0.5), Inches(1.4), Inches(7.3), Inches(3.5), size=19)
add_image(s, CAPTURES + "MVC.png", Inches(8), Inches(1.4), Inches(4.9), Inches(5.5))

# ── SLIDE 12 — Structure du projet ───────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Structure du projet — Atomic Design + Stimulus")
add_bullet_text(s, [
    "**Atoms** : bouton, lien, badge, input",
    "**Molecules** : carte produit, bouton panier, notation étoiles",
    "**Organisms** : header, footer, grille produits, panier live",
    "",
    "**14 controllers Stimulus** — interactions côté client :",
    "`carousel` · `cart-sidebar` · `favori` · `dropdown` · `confirm`",
    "`nav-toggle` · `image-zoom` · `star-rating` · `annulation`",
    "`password-toggle` · `flash-tooltip` · `consent-banner` · `csrf-protection` · `submit-once`",
], Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.5), size=20)
add_image(s, CAPTURES + "arborescenceProjet.png", Inches(8.2), Inches(1.4), Inches(2.2), Inches(5.5))
add_image(s, CAPTURES + "atomicDesignFichiers.png", Inches(10.6), Inches(1.4), Inches(2.3), Inches(5.5))

# ── SLIDE 13 — Code Atomic Design ────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Code : Atomic Design — Atom & Molécule")
add_textbox(s, "Atom — Button.html.twig", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """{% set variants = {
  primary: "bg-accent text-white hover:bg-primary hover:scale-105",
  ghost:   "bg-transparent text-text hover:bg-surface hover:scale-105",
  danger:  "bg-danger text-white hover:bg-danger-dark hover:scale-105",
  "ghost-danger": "text-danger border-2 border-danger hover:bg-danger",
} %}
<button type="{{ type }}"
  {{ attributes.defaults({
    class: base ~ ' ' ~ (variants[variant] ?? variants.primary)
  }) }}
>
  {% block content %}{% endblock %}
</button>""", Inches(0.5), Inches(1.7), Inches(6), Inches(2.8), lang='twig')

add_textbox(s, "Molécule — InputField.html.twig", Inches(0.5), Inches(4.6), Inches(6), Inches(0.35),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """<twig:Atoms:Label for="{{ id }}" label="{{ label }}"
  :required="required" />
<twig:Atoms:Input
  id="{{ id }}" name="{{ name }}" type="{{ type }}"
  aria-invalid="{{ isInvalid ? 'true' : 'false' }}"
  aria-describedby="{{ describedBy }}"
/>
{% if error %}
  <p id="{{ this.errorId }}" role="alert" class="text-danger">
    {{ error }}</p>
{% endif %}""", Inches(0.5), Inches(4.95), Inches(6), Inches(2.3), lang='twig')

add_textbox(s, "Utilisation :", Inches(7), Inches(1.35), Inches(6), Inches(0.35),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """<twig:Atoms:Button>Commander</twig:Atoms:Button>
<twig:Atoms:Button variant="ghost">Annuler</twig:Atoms:Button>

<twig:Molecules:InputField
  name="email"
  label="Email"
  type="email"
  :required="true"
/>""", Inches(7), Inches(1.7), Inches(5.9), Inches(2.3), lang='twig')

# ── SLIDE 14 — Code Carousel ─────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Code : Controller Stimulus — Carousel")
add_textbox(s, "Mode infini — clonage + saut silencieux", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35),
            size=14, bold=True, color=FRAMBOISE)
add_code_block(s, """this.items = [
  ...this.items.slice(this.items.length - this.offset)
               .map(i => i.cloneNode(true)),
  ...this.items,
  ...this.items.slice(0, this.offset).map(i => i.cloneNode(true)),
];""", Inches(0.5), Inches(1.7), Inches(6), Inches(1.8), lang='js')

add_textbox(s, "Protection double-clic — isAnimating", Inches(0.5), Inches(3.6), Inches(6), Inches(0.35),
            size=14, bold=True, color=FRAMBOISE)
add_code_block(s, """next() {
  if (this.isAnimating) return;
  this.isAnimating = true;
  this.gotoItem(this.currentItem + this.slidesToScroll);
  setTimeout(() => { this.isAnimating = false; },
    this.options.transitionDuration);
}""", Inches(0.5), Inches(3.95), Inches(6), Inches(1.9), lang='js')

add_textbox(s, "Accessibilité — attribut inert", Inches(7), Inches(1.35), Inches(6), Inches(0.35),
            size=14, bold=True, color=FRAMBOISE)
add_code_block(s, """setActiveItems() {
  this.items.forEach((item, i) => {
    const isActive =
      i >= this.currentItem &&
      i < this.currentItem + this.slidesVisible;
    item.inert = !isActive;
    item.setAttribute('aria-hidden', String(!isActive));
  });
}""", Inches(7), Inches(1.7), Inches(5.9), Inches(2.5), lang='js')

add_textbox(s, "Saut silencieux", Inches(7), Inches(4.3), Inches(6), Inches(0.35),
            size=14, bold=True, color=FRAMBOISE)
add_code_block(s, """onTransitionEnd() {
  if (this.currentItem <= this.offset) {
    this.setTransition(false);
    this.currentItem += this.items.length - 2 * this.offset;
    this.translate();
  }
}""", Inches(7), Inches(4.65), Inches(5.9), Inches(1.9), lang='js')

# ── SLIDE 15 — Base de données ───────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Base de données")
add_bullet_text(s, [
    "**6 entités :**",
    "`Utilisateur` · `Produit` · `Catégorie` · `Recette` · `Commande` · `Avis`",
    "**Table de liaison enrichie :** `CommandeProduit` (quantité + prix unitaire)",
    "**Panier :** session PHP — pas d'entité en base",
    "**Favoris :** `ManyToMany` sur `Utilisateur` — jointure automatique Doctrine",
    "",
    "**Points clés :**",
    "- Prix en `DECIMAL(8,2)` — pas de float (erreurs d'arrondi)",
    "- Prix dupliqué dans `CommandeProduit` — historique fiable",
    "- Enums PHP natifs (`StatutCommande`, `Difficulté`)",
], Inches(0.5), Inches(1.4), Inches(6.5), Inches(5), size=17)
add_image(s, CAPTURES + "DBdiagram.png", Inches(7.2), Inches(1.4), Inches(5.7), Inches(5.5))

# ── SLIDE 16 — Catalogue & Recettes ─────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Fonctionnalités : Catalogue & Recettes")
add_bullet_text(s, [
    "**Catalogue produits**",
    "- Filtres par catégorie, tri par prix/note",
    "- Turbo Frames — navigation sans rechargement",
    "- Zoom image au clic",
    "- Système de favoris (AJAX)",
    "",
    "**Section recettes**",
    "- Même organism que le catalogue — recettes gratuites, même expérience",
], Inches(0.5), Inches(1.4), Inches(6), Inches(5), size=18)
add_image(s, CAPTURES + "grilleProduitDesktop.png", Inches(6.7), Inches(1.4), Inches(3), Inches(2.7))
add_image(s, CAPTURES + "ReccetteGrid.png", Inches(10), Inches(1.4), Inches(3), Inches(2.7))

# ── SLIDE 17 — Panier ────────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Fonctionnalités : Panier")
add_bullet_text(s, [
    "**Sidebar panier persistante**",
    "- Ouverture/fermeture via Stimulus",
    "- Mise à jour en temps réel (Live Component)",
    "- Ajout / retrait / suppression de produits",
    "- Confirmation avant vidage (dialog natif HTML)",
], Inches(0.5), Inches(1.4), Inches(6.5), Inches(4), size=19)
add_image(s, CAPTURES + "modalPanier.png", Inches(7), Inches(1.4), Inches(2.9), Inches(4.5))
add_image(s, CAPTURES + "ConfirmationViderPanier.png", Inches(10.1), Inches(1.4), Inches(2.9), Inches(4.5))

# ── SLIDE 18 — Commande & Paiement ───────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Fonctionnalités : Commande & Paiement")
add_bullet_text(s, [
    "**Tunnel de commande**",
    "1. Récapitulatif panier",
    "2. Redirection vers Stripe Checkout",
    "3. Webhook Stripe → mise à jour statut commande",
    "4. Email de confirmation + facture PDF",
], Inches(0.5), Inches(1.4), Inches(6.5), Inches(4), size=19)
add_image(s, CAPTURES + "StripeCheckout.png", Inches(7.2), Inches(1.4), Inches(5.7), Inches(5.5))

# ── SLIDE 19 — Emails ────────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Fonctionnalités : Emails transactionnels")
add_table(s,
    ["Déclencheur", "Email envoyé"],
    [
        ["Inscription",          "Lien de vérification"],
        ["Compte vérifié",       "Email de bienvenue"],
        ["Commande confirmée",   "Confirmation + facture PDF"],
        ["Commande annulée",     "Notification d'annulation"],
    ],
    Inches(0.5), Inches(1.5), Inches(12.3)
)
add_textbox(s, "Service : Resend API — domaine samydessert.fr acheté et vérifié via DNS",
            Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
            size=13, italic=True, color=FRAMBOISE)

# ── SLIDE 20 — Administration ────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Fonctionnalités : Administration")
add_bullet_text(s, [
    "**Interface EasyAdmin**",
    "- Gestion produits, recettes, commandes, utilisateurs, avis",
    "- Upload d'images via VichUploader (stockage local)",
    "- Accès réservé `ROLE_ADMIN`",
    "",
    "**Système d'avis**",
    "- Note 1-5 étoiles, note moyenne calculée dynamiquement",
    "- Un seul avis par utilisateur par produit",
], Inches(0.5), Inches(1.4), Inches(6), Inches(5), size=18)
add_image(s, CAPTURES + "adminListeProduits.png", Inches(6.7), Inches(1.4), Inches(3), Inches(2.7))
add_image(s, CAPTURES + "adminEditionProduit.png", Inches(10), Inches(1.4), Inches(3), Inches(2.7))

# ── SLIDE 21 — Soins UX ──────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Soins UX")
add_bullet_text(s, [
    "- Messages d'erreur vagues (connexion) — recommandation OWASP",
    "- Bouton panier — Live Component : spinner pendant l'action, re-render automatique",
], Inches(0.5), Inches(1.4), Inches(6.5), Inches(1.2), size=18)
add_image(s, CAPTURES + "ConnexionErreur.png", Inches(0.5), Inches(2.7), Inches(3.5), Inches(4.5))
add_textbox(s, "BoutonPanier.php", Inches(4.2), Inches(1.38), Inches(4.5), Inches(0.4),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """#[AsLiveComponent]
final class BoutonPanier
{
    #[LiveProp]
    public int $produitId = 0;

    public function getQuantite(): int
    {
        return $this->panier
            ->getQuantitePourProduit($this->produitId);
    }

    #[LiveAction]
    public function ajouter(): void
    {
        $this->panier->ajouter($this->produitId);
        $this->emit('panierUpdated');
    }
}""", Inches(4.2), Inches(1.75), Inches(4.5), Inches(4.5), lang='php')
add_textbox(s, "Twig — spinner pendant le chargement", Inches(9.0), Inches(1.38), Inches(4.0), Inches(0.4),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """<twig:Atoms:Icon
  name="cart-plus"
  data-loading="hide"
/>
<twig:Atoms:Spinner
  data-loading="show"
/>""", Inches(9.0), Inches(1.75), Inches(4.0), Inches(2.5), lang='twig')

# ── SLIDE 22 — Sécurité ──────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Sécurité")
add_table(s,
    ["Menace", "Protection"],
    [
        ["CSRF",                "Token Symfony sur tous les formulaires"],
        ["XSS",                 "Échappement automatique Twig"],
        ["Injection SQL",       "Doctrine ORM (requêtes préparées)"],
        ["Accès non autorisé",  "ROLE_ADMIN / ROLE_USER + access_control"],
        ["Mots de passe",       "Hachage automatique Symfony (auto) — Sodium/bcrypt"],
        ["Données de paiement", "Jamais stockées — délégation à Stripe"],
    ],
    Inches(0.5), Inches(1.4), Inches(9)
)
add_image(s, CAPTURES + "securiteYaml.png", Inches(9.8), Inches(1.4), Inches(3.1), Inches(5.5))

# ── SLIDE 23 — Code Webhook ───────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Code : Webhook Stripe & idempotence")
add_textbox(s, "1 — Vérification de signature", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """try {
    $event = Webhook::constructEvent(
        $payload, $sigHeader, $webhookSecret
    );
} catch (\\UnexpectedValueException $e) {
    return new Response('Invalid payload', 400);
} catch (SignatureVerificationException $e) {
    return new Response('Invalid signature', 400);
}""", Inches(0.5), Inches(1.7), Inches(6), Inches(2.6), lang='php')

add_textbox(s, "2 — Idempotence : vérifier le statut avant d'agir",
            Inches(6.9), Inches(1.35), Inches(6), Inches(0.35),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """if ($event->type === 'checkout.session.completed') {
    $commande = $this->commandeRepo
        ->findOneByStripeSessionId($stripeSession->id);

    if ($commande &&
        $commande->getStatut() === StatutCommande::EnAttente) {
        $commande->setStatut(StatutCommande::Confirmee);
        $this->em->flush();
        $this->mailer->envoyerConfirmationCommande($commande);
    }
}
return new Response('OK', 200);""", Inches(6.9), Inches(1.7), Inches(6), Inches(3.2), lang='php')

# ── SLIDE 24 — Code Favoris ───────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Code : Favoris AJAX — Stimulus ↔ PHP")
add_textbox(s, "PHP — une route pour l'ajout ET le retrait", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """#[Route('/favori/{type}/{id}', methods: ['POST'])]
public function toggle(string $type, int $id,
    EntityManagerInterface $em): JsonResponse
{
    if ($user->getProduitsFavoris()->contains($entity)) {
        $user->removeProduitFavori($entity);
        $favori = false;
    } else {
        $user->addProduitFavori($entity);
        $favori = true;
    }
    $em->flush();
    return $this->json(['favori' => $favori]);
}""", Inches(0.5), Inches(1.7), Inches(6.2), Inches(3.5), lang='php')

add_textbox(s, "JS — fetch + communication entre controllers Stimulus",
            Inches(6.9), Inches(1.35), Inches(6), Inches(0.35),
            size=13, bold=True, color=FRAMBOISE)
add_code_block(s, """async toggle(e) {
    e.preventDefault();
    const response = await fetch(this.urlValue, {
        method: 'POST',
        headers: { 'X-Requested-With': 'XMLHttpRequest' }
    });
    if (response.status === 401) {
        this.element
            .closest('[data-controller~="flash-tooltip"]')
            ?.__stimulus_flash_tooltip?.show();
        return;
    }
    const data = await response.json();
    this.activeValue = data.favori;
}""", Inches(6.9), Inches(1.7), Inches(6), Inches(3.5), lang='js')

# ── SLIDE 25 — Code PanierService ────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Code : PanierService — Service & accès aux données")
add_textbox(s, "Le Controller délègue — toute la logique est dans le Service",
            Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.35),
            size=14, bold=True, color=FRAMBOISE)
add_code_block(s, """public function getLignes(): array
{
    $panier   = $this->getPanierSession();
    $produits = $this->produitRepository->findBy(['id' => array_keys($panier)]);
    foreach ($produits as $produit) {
        $lignes[] = ['produit' => $produit, 'quantite' => $panier[$produit->getId()]];
    }
    return $lignes;
}

public function ajouter(int $produitId): void
{
    $panier = $this->getPanierSession();
    $panier[$produitId] = ($panier[$produitId] ?? 0) + 1;
    $this->savePanier($panier);
}

public function getTotal(): float
{
    return array_reduce($this->getLignes(), fn($total, $ligne) =>
        $total + $ligne['produit']->getPrix() * $ligne['quantite'], 0.0
    );
}""", Inches(0.5), Inches(1.75), Inches(12.3), Inches(5), lang='php')

# ── SLIDE 26 — Tests ─────────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Tests — PHPUnit 12 : 74 tests, 0 erreurs")
add_table(s,
    ["Type", "Nombre", "Ce qui est testé"],
    [
        ["Unitaires",   "26", "Logique métier isolée — panier, commande, emails, accès compte"],
        ["Fonctionnels","48", "Requêtes HTTP réelles — routes, formulaires, AJAX"],
    ],
    Inches(0.5), Inches(1.5), Inches(12.3)
)
add_textbox(s, "Base de test isolée (samydessert_test), réinitialisée entre chaque test",
            Inches(0.5), Inches(3), Inches(12.3), Inches(0.4), size=13, italic=True, color=CHOCOLAT)
add_textbox(s, "Exemple : test corrigé — token CSRF manquant",
            Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.4), size=15, bold=True, color=FRAMBOISE)
add_code_block(s, """// Avant : aucun token → redirection vers /connexion
$client->request('POST', '/contact', [...]);

// Après : on récupère le vrai token depuis le formulaire
$client->request('GET', '/contact');
$token = $client->getCrawler()->filter('input[name="_token"]')->attr('value');
$client->request('POST', '/contact', [..., '_token' => $token]);
$this->assertResponseRedirects('/contact');""",
    Inches(0.5), Inches(3.9), Inches(9.5), Inches(2.6), lang='php')
add_image(s, CAPTURES + "TestPHPunit.png", Inches(10.2), Inches(3.9), Inches(2.7), Inches(2.6))

# ── SLIDE 27 — Déploiement ───────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Déploiement")
add_bullet_text(s, [
    "**Environnement de développement**",
    "- Docker Compose — 6 services (nginx, php-fpm, mysql, adminer, init, assets)",
    "- Pas d'installation locale requise",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.8), size=19)
add_rect(s, Inches(0.5), Inches(3.1), Inches(12.3), Pt(1), fill=FRAMBOISE)
add_bullet_text(s, [
    "**Production — Railway**",
    "- Déploiement automatique sur push `main`",
    "- 1 container Docker (Dockerfile à la racine)",
    "- Variables sensibles dans Railway (jamais commitées)",
    "- URL : https://samydessert-production.up.railway.app",
], Inches(0.5), Inches(3.2), Inches(12.3), Inches(3), size=19)

# ── SLIDE 28 — Git ───────────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Gestion de version avec Git")
add_table(s,
    ["Branche", "Rôle"],
    [
        ["`dev`",  "Développement — isolation des changements en cours"],
        ["`main`", "Code stable → déploiement Railway automatique"],
    ],
    Inches(0.5), Inches(1.5), Inches(8)
)
add_bullet_text(s, [
    "**Convention de commits — Conventional Commits :**",
    "- `feat:` nouvelle fonctionnalité",
    "- `fix:` correction de bug",
    "- `docs:` documentation",
    "- `refactor:` refactoring",
], Inches(0.5), Inches(3.5), Inches(8), Inches(3.5), size=18)

# ── SLIDE 29 — Veille ────────────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Veille technologique")
add_bullet_text(s, [
    "**Comment je me tiens informé**",
    "- Stack Overflow, GitHub Issues — résolution de problèmes concrets",
    "- YouTube (Grafikart, Symfony Online Conference)",
    "- **IA (Claude, ChatGPT)** — compréhension rapide, vérification des réponses",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(2), size=18)
add_table(s,
    ["Technologie", "Nouveauté"],
    [
        ["Symfony 7.4",     "Live Components, UX Twig Components"],
        ["Tailwind CSS v4", "Tokens CSS natifs, pas de tailwind.config.js"],
        ["PHP 8.3",         "Enums natifs, readonly properties"],
        ["Resend API",      "Alternative moderne à SMTP pour les emails en prod"],
        ["AssetMapper",     "Remplace Webpack Encore — zéro configuration"],
    ],
    Inches(0.5), Inches(3.5), Inches(12.3)
)

# ── SLIDE 30 — Difficultés ───────────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Difficultés rencontrées")
add_table(s,
    ["Difficulté", "Solution"],
    [
        ["Webhook Stripe en local",          "Stripe CLI (stripe listen)"],
        ["Emails bloqués depuis Railway",    "Migration SMTP → Resend API (HTTP)"],
        ["Tailwind watch sur Docker Windows","Script de recompilation toutes les 3s"],
        ["Turbo Frames & rechargement",      "Identification précise des frames Turbo"],
        ["VichUploader + AssetMapper",       "Configuration manuelle du mapping"],
    ],
    Inches(0.5), Inches(1.5), Inches(12.3)
)

# ── SLIDE 31 — Améliorations futures ─────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Améliorations futures")
add_bullet_text(s, [
    "**Fonctionnalités**",
    "- Liaison produit ↔ recette (OneToOne — créer la recette à la création du produit)",
    "- Ingrédients dans les recettes",
    "- Modération des avis avant publication",
    "- Variations de produits (taille, parfum)",
    "- Recettes proposées par les utilisateurs",
], Inches(0.5), Inches(1.4), Inches(6), Inches(4), size=18)
add_bullet_text(s, [
    "**Technique**",
    "- Remplacer `php -S` par Nginx + PHP-FPM en prod",
    "- Token de vérification email avec expiration",
    "- Token CSRF sur les actions AJAX (favoris)",
    "- Internationalisation (symfony/translation)",
], Inches(6.9), Inches(1.4), Inches(6), Inches(4), size=18)

# ── SLIDE 32 — Démonstration live ────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Démonstration live")
add_textbox(s, "https://samydessert-production.up.railway.app",
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
            size=18, bold=True, color=FRAMBOISE, align=PP_ALIGN.CENTER)
add_bullet_text(s, [
    "1. Accueil → catalogue → fiche produit",
    "2. Ajout au panier → commande → paiement Stripe",
    "3. Interface admin",
    "",
    "Comptes de démonstration :",
    "- Admin : samy@admin.com / password123",
    "- Utilisateur : test@test.com / password123",
], Inches(1), Inches(2.2), Inches(11.3), Inches(4.5), size=19)

# ── SLIDE 33 — Bilan personnel ───────────────────────────────────────────────
s = blank_slide(prs)
add_title_bar(s, "Bilan personnel")
add_bullet_text(s, [
    "**Ce que j'ai appris**",
    "- Structurer un projet Symfony de A à Z",
    "- Travailler avec Docker en environnement de dev",
    "- Intégrer des services tiers (Stripe, Resend, VichUploader)",
    "- Utiliser l'IA comme outil d'apprentissage — poser les bonnes questions, vérifier les réponses",
    "- Déployer en production sur Railway",
    "- Concevoir une identité visuelle cohérente avec Affinity Designer et Figma",
    "- Appliquer les principes d'accessibilité web (WCAG, Luciole, focus clavier)",
], Inches(0.5), Inches(1.4), Inches(8), Inches(5), size=20)
add_rect(s, Inches(8.7), Inches(1.4), Pt(2), Inches(5.5), fill=FRAMBOISE)
add_textbox(s,
    "Voir un projet passer d'une idée à un produit réel, en ligne, utilisable — et comprendre chaque brique qui le compose.",
    Inches(9), Inches(2.5), Inches(4), Inches(2.5),
    size=16, italic=True, color=FRAMBOISE, align=PP_ALIGN.CENTER)

# ── SLIDE 34 — Questions ─────────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, W, H, fill=CHOCOLAT)
add_rect(s, 0, Inches(3), W, Inches(1.5), fill=FRAMBOISE)
add_textbox(s, "Merci !", Inches(1), Inches(3.05), Inches(11), Inches(1.2),
            size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "SamyDessert — Boutique en ligne de pâtisseries artisanales",
            Inches(1), Inches(5), Inches(11), Inches(0.5),
            size=16, color=CREME, align=PP_ALIGN.CENTER)
add_textbox(s, "Prêt pour vos questions",
            Inches(1), Inches(5.6), Inches(11), Inches(0.5),
            size=14, italic=True, color=GREY_LIGHT, align=PP_ALIGN.CENTER)

# ── SAUVEGARDE ───────────────────────────────────────────────────────────────
output = r"c:\Travail\samyDessert\Soutenance\diaporama\SamyDessert_Soutenance.pptx"
prs.save(output)
print(f"Fichier généré : {output}")
